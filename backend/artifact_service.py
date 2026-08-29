"""Generate polished, editable, user-downloadable Ask Crump artifacts.

The model writes the source content. This deterministic finishing layer applies
format-aware typography, safe spreadsheet behavior, restrained branding,
document metadata, and layouts that remain editable in Office applications.
"""
from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
import re
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor as DocxRGBColor
from openpyxl import Workbook
from openpyxl.formatting.rule import DataBarRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.worksheet.table import Table as XlsxTable, TableStyleInfo
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle,
)

from .file_service import FileService


VALID_FORMATS = {'docx', 'pdf', 'pptx', 'xlsx', 'md', 'txt'}
MIME = {
    'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
    'pdf': 'application/pdf',
    'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
    'md': 'text/markdown',
    'txt': 'text/plain',
}
INK = '171B24'
NAVY = '202532'
GOLD = 'C9A95E'
CREAM = 'F6F1E7'
MUTED = '677080'
LINE = 'D9DDE4'
PAPER = 'FAF9F6'
SLATE = '405366'
SOFT = 'EEF1F4'
TEAL = '1F6A67'
COPPER = 'B85C3B'
SKY = '8FB7C5'
MIST = 'E8F0EF'
WHITE = 'FFFFFF'


@dataclass(slots=True)
class ArtifactService:
    files: FileService

    ACTION_PATTERN = re.compile(
        r"\b(create|make|generate|export|deliver|build|produce|turn|convert|save|"
        r"put|write|draft|compose|prepare|format|package|send|provide|give|download)\b", re.I,
    )
    FOLLOW_UP_DELIVERY_PATTERN = re.compile(
        r"\b(export|deliver|download|package|attach)\b|"
        r"\b(make|turn|put|convert|send|provide|give)\b.{0,80}"
        r"\b(downloadable|file|attachment|document|doc)\b",
        re.I,
    )
    FOLLOW_UP_REFERENCE_PATTERN = re.compile(
        r"\b(this|that|it|the above|what you wrote|your draft|the draft|"
        r"downloadable|file|attachment|document|doc)\b",
        re.I,
    )
    NON_DOCUMENT_CREATION_PATTERN = re.compile(
        r"\b(image|picture|photo|photograph|artwork|illustration|logo|poster|cover|"
        r"video|movie|film|clip|animation|scene|book|novel|memoir|screenplay)\b",
        re.I,
    )
    LONG_FORM_PATTERN = re.compile(
        r"\b(novel|book|memoir|manuscript|screenplay|dissertation|thesis)\b", re.I,
    )
    SHORT_FORM_PATTERN = re.compile(
        r"\b(summary|synopsis|outline|review|blurb|proposal|query letter|book report)\b", re.I,
    )
    RESUME_PATTERN = re.compile(
        r"\b(resume|r[ée]sum[ée]|curriculum vitae|\bcv\b|cover letter|job application)\b", re.I,
    )
    ACADEMIC_PATTERN = re.compile(
        r"\b(college|university|academic|essay|research paper|term paper|literature review|"
        r"dissertation|thesis|mla|apa|chicago style|works cited|bibliography)\b", re.I,
    )

    @staticmethod
    def normalize_format(value: Any) -> str | None:
        raw = str(value or '').strip().lower().lstrip('.')
        aliases = {'word': 'docx', 'powerpoint': 'pptx', 'excel': 'xlsx', 'markdown': 'md', 'text': 'txt'}
        raw = aliases.get(raw, raw)
        return raw if raw in VALID_FORMATS else None

    @classmethod
    def is_long_form_request(cls, message: str) -> bool:
        text = str(message or '').strip()
        if not text or not cls.ACTION_PATTERN.search(text) or not cls.LONG_FORM_PATTERN.search(text):
            return False
        if cls.SHORT_FORM_PATTERN.search(text):
            return False
        explicit_length = bool(re.search(
            r"\b(full[ -]?length|complete|entire|book[ -]?length|from start to finish)\b|"
            r"\b\d{2,3}(?:,\d{3})+\s*words?\b", text, re.I,
        ))
        writing_verb = bool(re.search(r"\b(write|draft|compose|author|create|produce|build)\b", text, re.I))
        substantial_kind = bool(re.search(
            r"\b(novel|memoir|manuscript|screenplay|dissertation|thesis)\b", text, re.I,
        ))
        return explicit_length or (writing_verb and substantial_kind) or bool(
            writing_verb and re.search(r"\bbook\b", text, re.I)
        )

    @classmethod
    def _mentioned_format(cls, value: Any, *, allow_generic: bool = False) -> str | None:
        text = str(value or '').lower().strip()
        if not text:
            return None
        formats = [
            (r'\b(powerpoint|pptx|slide deck|presentation)\b', 'pptx'),
            (r'\b(excel|xlsx|spreadsheet|workbook)\b', 'xlsx'),
            (r'\bdocx\b|\bmicrosoft\s+word\b|\bword\s+(?:document(?:ed)?|doc|file|manuscript)\b|\bdoc\b', 'docx'),
            (r'\bpdf\b', 'pdf'),
            (r'\bmarkdown|\.md\b', 'md'),
            (r'\btext file|\.txt\b', 'txt'),
        ]
        for pattern, fmt in formats:
            if re.search(pattern, text):
                return fmt
        if allow_generic and re.search(
            r'\b(document|manuscript|report|letter|resume|r[ée]sum[ée]|essay|paper|proposal|cv)\b',
            text,
        ):
            return 'docx'
        return None

    @classmethod
    def detect_request(
        cls,
        message: str,
        explicit: Any = None,
        history: Any = None,
    ) -> str | None:
        selected = cls.normalize_format(explicit)
        if selected:
            return selected
        text = str(message or '').lower().strip()
        if not cls.ACTION_PATTERN.search(text):
            return None
        mentioned = cls._mentioned_format(text)
        if mentioned:
            return mentioned
        if re.search(
            r'\b(document|manuscript|report|letter|resume|r[ée]sum[ée]|essay|paper|proposal|cv)\b',
            text,
        ) and re.search(r'\b(file|downloadable|download|attachment|send|deliver|export|document)\b', text):
            return 'docx'
        if not (
            cls.FOLLOW_UP_DELIVERY_PATTERN.search(text)
            and cls.FOLLOW_UP_REFERENCE_PATTERN.search(text)
            and isinstance(history, list)
        ):
            return None
        for item in reversed(history[-12:]):
            if not isinstance(item, dict) or str(item.get('role') or '').lower() != 'user':
                continue
            content = item.get('content')
            mentioned = cls._mentioned_format(content, allow_generic=True)
            if mentioned:
                return mentioned
            if cls.NON_DOCUMENT_CREATION_PATTERN.search(str(content or '')):
                return None
        return None

    @classmethod
    def normalize_purpose(cls, value: Any = None) -> str | None:
        purpose = str(value or '').strip().lower()
        return purpose if purpose in {'resume'} else None

    @classmethod
    def profile_for(
        cls,
        markdown: str = '',
        brief: str = '',
        format_name: str = '',
        purpose: Any = None,
    ) -> str:
        if cls.normalize_purpose(purpose) == 'resume':
            return 'resume'
        evidence = f'{brief}\n{markdown[:8000]}'.strip()
        if cls.RESUME_PATTERN.search(evidence):
            return 'resume'
        if cls.ACADEMIC_PATTERN.search(evidence):
            return 'academic'
        if cls.normalize_format(format_name) == 'pptx':
            return 'presentation'
        if cls.normalize_format(format_name) == 'xlsx':
            return 'spreadsheet'
        return 'business'

    @classmethod
    def creation_guidance(cls, format_name: str, brief: str = '', purpose: Any = None) -> str:
        """Return a model-facing contract for the requested deliverable."""
        fmt = cls.normalize_format(format_name) or 'docx'
        profile = cls.profile_for('', brief, fmt, purpose)
        shared = (
            'Return only the finished source content—no production notes, promises, or commentary about '
            'creating a file. Follow every supplied fact and requirement. Never invent citations, employers, '
            'degrees, dates, metrics, quotations, research findings, or financial inputs. When a required fact '
            'is unavailable, use a clearly labeled placeholder or state the limitation instead of fabricating it.'
        )
        if profile == 'resume':
            return (
                f'{shared} Build an ATS-friendly résumé with a concise professional summary, relevant skills, '
                'reverse-chronological experience, and accomplishment bullets. Quantify impact only when the '
                'user supplied the number. Do not use tables, columns, icons, headshots, ratings, or references.'
            )
        if profile == 'academic':
            return (
                f'{shared} Build a coherent academic draft with a defensible thesis, logical sections, evidence '
                'and analysis, an appropriate counterpoint, and a substantive conclusion. Cite only sources '
                'present in supplied context or verified research results. Match the requested citation style; '
                'include a references section only when grounded source details are available.'
            )
        if fmt == 'pptx':
            return (
                f'{shared} Write a presentation outline in Markdown. Use one clear takeaway headline per slide, '
                'short supporting points, and no more than five points per slide. Use compact Markdown tables '
                'for grounded comparisons or numeric data that should become an editable chart. Prefer a '
                'narrative arc: context, insight, recommendation, execution, next step. Put grounded sources '
                'in a final Sources section.'
            )
        if fmt == 'xlsx':
            return (
                f'{shared} Provide structured Markdown tables with explicit assumptions, inputs, units, and '
                'outputs. Put formulas in cells beginning with = and keep formulas limited to normal cell '
                'references and standard spreadsheet functions. Never invent business or financial data.'
            )
        return (
            f'{shared} Use a clear title, executive-quality hierarchy, concise paragraphs, descriptive headings, '
            'useful lists, and tables where comparison or structured data benefits from them.'
        )

    @staticmethod
    def title_from(text: str, fallback: str = 'Ask Crump') -> str:
        for line in str(text or '').splitlines():
            clean = re.sub(r'^#{1,6}\s*', '', line).strip()
            if clean and not clean.startswith('|'):
                return clean[:160]
        return fallback

    @staticmethod
    def safe_stem(title: str) -> str:
        stem = re.sub(r'[^A-Za-z0-9 -]+', '', title).strip()
        stem = re.sub(r'\s+', '_', stem)[:60]
        return stem or 'Ask_Crump_Document'

    @staticmethod
    def _table_separator(line: str) -> bool:
        cells = [cell.strip() for cell in str(line).strip().strip('|').split('|')]
        return bool(cells) and all(re.fullmatch(r':?-{3,}:?', cell or '') for cell in cells)

    @classmethod
    def blocks(cls, markdown: str) -> list[tuple[str, Any]]:
        """Parse the conservative Markdown subset used across exported formats."""
        source = str(markdown or '').replace('\r\n', '\n').split('\n')
        output: list[tuple[str, Any]] = []
        paragraph: list[str] = []
        code: list[str] = []
        in_code = False
        index = 0

        def flush_para() -> None:
            if paragraph:
                output.append(('p', ' '.join(part.strip() for part in paragraph if part.strip())))
                paragraph.clear()

        while index < len(source):
            raw = source[index]
            stripped = raw.strip()
            if stripped.startswith('```'):
                flush_para()
                if in_code:
                    output.append(('code', '\n'.join(code)))
                    code.clear()
                in_code = not in_code
                index += 1
                continue
            if in_code:
                code.append(raw)
                index += 1
                continue
            if not stripped:
                flush_para()
                index += 1
                continue
            if '|' in stripped and index + 1 < len(source) and cls._table_separator(source[index + 1]):
                flush_para()
                rows = [[cell.strip() for cell in stripped.strip('|').split('|')]]
                index += 2
                while index < len(source) and '|' in source[index] and source[index].strip():
                    rows.append([cell.strip() for cell in source[index].strip().strip('|').split('|')])
                    index += 1
                output.append(('table', rows))
                continue
            heading = re.match(r'^(#{1,4})\s+(.+)$', stripped)
            if heading:
                flush_para()
                output.append((f'h{len(heading.group(1))}', heading.group(2).strip()))
            else:
                bullet = re.match(r'^[-*+]\s+(.+)$', stripped)
                ordered = re.match(r'^\d+[.)]\s+(.+)$', stripped)
                quote = re.match(r'^>\s?(.+)$', stripped)
                if bullet:
                    flush_para(); output.append(('li', bullet.group(1).strip()))
                elif ordered:
                    flush_para(); output.append(('ol', ordered.group(1).strip()))
                elif quote:
                    flush_para(); output.append(('quote', quote.group(1).strip()))
                elif re.fullmatch(r'[-*_]{3,}', stripped):
                    flush_para(); output.append(('hr', ''))
                else:
                    paragraph.append(stripped)
            index += 1
        flush_para()
        if code:
            output.append(('code', '\n'.join(code)))
        return output

    @staticmethod
    def _set_cell_shading(cell: Any, color: str) -> None:
        props = cell._tc.get_or_add_tcPr()
        shading = props.find(qn('w:shd'))
        if shading is None:
            shading = OxmlElement('w:shd'); props.append(shading)
        shading.set(qn('w:fill'), color)

    @staticmethod
    def _set_cell_margins(cell: Any, top: int = 90, start: int = 100, bottom: int = 90, end: int = 100) -> None:
        props = cell._tc.get_or_add_tcPr()
        margins = props.first_child_found_in('w:tcMar')
        if margins is None:
            margins = OxmlElement('w:tcMar'); props.append(margins)
        for name, value in {'top': top, 'start': start, 'bottom': bottom, 'end': end}.items():
            node = margins.find(qn(f'w:{name}'))
            if node is None:
                node = OxmlElement(f'w:{name}'); margins.append(node)
            node.set(qn('w:w'), str(value)); node.set(qn('w:type'), 'dxa')

    @staticmethod
    def _add_page_number(paragraph: Any) -> None:
        paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = paragraph.add_run()
        begin = OxmlElement('w:fldChar'); begin.set(qn('w:fldCharType'), 'begin')
        instruction = OxmlElement('w:instrText'); instruction.set(qn('xml:space'), 'preserve'); instruction.text = ' PAGE '
        end = OxmlElement('w:fldChar'); end.set(qn('w:fldCharType'), 'end')
        run._r.extend([begin, instruction, end])

    @staticmethod
    def _clean_inline(text: str) -> str:
        value = re.sub(r'!\[([^]]*)\]\([^)]+\)', r'\1', str(text or ''))
        value = re.sub(r'\[([^]]+)\]\([^)]+\)', r'\1', value)
        return re.sub(r'[`*_~]+', '', value).strip()

    @classmethod
    def _same_title(cls, candidate: str, resolved_title: str) -> bool:
        left = re.sub(r'\s+', ' ', cls._clean_inline(candidate)).strip().casefold()
        right = re.sub(r'\s+', ' ', cls._clean_inline(resolved_title)).strip().casefold()
        if not left or not right:
            return False
        return left == right or (min(len(left), len(right)) >= 30 and (left.startswith(right) or right.startswith(left)))

    @classmethod
    def _add_docx_inline(cls, paragraph: Any, text: str) -> None:
        pattern = re.compile(
            r'(\*\*[^*]+\*\*|__[^_]+__|(?<!\*)\*[^*]+\*(?!\*)|(?<!_)_[^_]+_(?!_)|`[^`]+`|\[[^]]+\]\([^)]+\))'
        )
        cursor = 0
        for match in pattern.finditer(str(text or '')):
            if match.start() > cursor:
                paragraph.add_run(text[cursor:match.start()])
            token = match.group(0); run = paragraph.add_run()
            if token.startswith(('**', '__')):
                run.text = token[2:-2]; run.bold = True
            elif token.startswith(('*', '_')):
                run.text = token[1:-1]; run.italic = True
            elif token.startswith('`'):
                run.text = token[1:-1]; run.font.name = 'Consolas'; run.font.size = Pt(9)
            else:
                link = re.match(r'\[([^]]+)\]\(([^)]+)\)', token)
                run.text = link.group(1) if link else token
                run.font.color.rgb = DocxRGBColor(33, 88, 148); run.underline = True
            cursor = match.end()
        if cursor < len(text):
            paragraph.add_run(text[cursor:])

    @staticmethod
    def _shade_paragraph(paragraph: Any, color: str) -> None:
        shading = OxmlElement('w:shd'); shading.set(qn('w:fill'), color)
        paragraph._p.get_or_add_pPr().append(shading)

    @staticmethod
    def _set_docx_font(font: Any, style_element: Any, name: str) -> None:
        """Set every Word font slot so Office themes cannot silently replace it."""
        font.name = name
        r_pr = style_element.get_or_add_rPr()
        r_fonts = r_pr.rFonts
        if r_fonts is None:
            r_fonts = OxmlElement('w:rFonts')
            r_pr.insert(0, r_fonts)
        for theme_slot in ('asciiTheme', 'hAnsiTheme', 'eastAsiaTheme', 'cstheme'):
            r_fonts.attrib.pop(qn(f'w:{theme_slot}'), None)
        for slot in ('ascii', 'hAnsi', 'eastAsia', 'cs'):
            r_fonts.set(qn(f'w:{slot}'), name)

    @staticmethod
    def _set_docx_size(font: Any, style_element: Any, size: float) -> None:
        font.size = Pt(size)
        r_pr = style_element.get_or_add_rPr()
        for element_name in ('sz', 'szCs'):
            node = r_pr.find(qn(f'w:{element_name}'))
            if node is None:
                node = OxmlElement(f'w:{element_name}')
                r_pr.append(node)
            node.set(qn('w:val'), str(round(size * 2)))

    @classmethod
    def _column_ratios(cls, rows: list[list[Any]], width: int) -> list[float]:
        """Infer practical table widths from the content instead of forcing equal columns."""
        weights: list[float] = []
        for column in range(width):
            lengths = [len(cls._clean_inline(str(row[column] if column < len(row) else ''))) for row in rows[:60]]
            longest = max(lengths, default=8)
            average = sum(lengths) / max(1, len(lengths))
            weights.append(min(42.0, max(8.0, average * 0.65 + longest * 0.35)))
        total = sum(weights) or 1.0
        return [weight / total for weight in weights]

    def docx(self, markdown: str, *, profile: str | None = None, title: str | None = None) -> bytes:
        resolved_title = (title or self.title_from(markdown)).strip()[:160]
        resolved_profile = profile or self.profile_for(markdown)
        doc = Document(); section = doc.sections[0]
        if resolved_profile == 'academic':
            margins = (1.0, 1.0, 1.0, 1.0); body_font, body_size, spacing = 'Times New Roman', 12, 2.0
            display_font = 'Times New Roman'
        elif resolved_profile == 'resume':
            margins = (0.62, 0.67, 0.62, 0.67); body_font, body_size, spacing = 'Aptos', 10, 1.03
            display_font = 'Aptos'
        else:
            margins = (0.82, 0.82, 0.78, 0.82); body_font, body_size, spacing = 'Aptos', 10.5, 1.14
            display_font = 'Aptos'
        section.top_margin, section.right_margin = Inches(margins[0]), Inches(margins[1])
        section.bottom_margin, section.left_margin = Inches(margins[2]), Inches(margins[3])

        normal = doc.styles['Normal']; self._set_docx_font(normal.font, normal.element, body_font); self._set_docx_size(normal.font, normal.element, body_size)
        normal.font.color.rgb = DocxRGBColor.from_string(INK)
        normal.paragraph_format.line_spacing = spacing
        normal.paragraph_format.space_after = Pt(0 if resolved_profile == 'academic' else 6)
        normal.paragraph_format.widow_control = True
        if resolved_profile == 'academic':
            normal.paragraph_format.first_line_indent = Inches(0.5)
        sizes = {'academic': (12, 12, 12, 12), 'resume': (22, 10.5, 10.5, 10), 'business': (24, 15, 12.5, 10.5)}[
            resolved_profile if resolved_profile in {'academic', 'resume'} else 'business'
        ]
        for name, size in zip(('Title', 'Heading 1', 'Heading 2', 'Heading 3'), sizes):
            style = doc.styles[name]; self._set_docx_font(style.font, style.element, display_font); self._set_docx_size(style.font, style.element, size)
            style.font.bold = True
            style.font.italic = resolved_profile == 'academic' and name == 'Heading 3'
            style.font.color.rgb = DocxRGBColor.from_string(INK if resolved_profile in {'academic', 'resume'} else SLATE)
            style.paragraph_format.keep_with_next = True; style.paragraph_format.keep_together = True
            if resolved_profile == 'academic':
                style.paragraph_format.line_spacing = 2.0
                style.paragraph_format.space_before = Pt(0)
                style.paragraph_format.space_after = Pt(0)
            elif resolved_profile == 'resume':
                style.paragraph_format.space_before = Pt(8 if name != 'Title' else 0)
                style.paragraph_format.space_after = Pt(3 if name != 'Title' else 4)
            else:
                style.paragraph_format.space_before = Pt(13 if name != 'Title' else 0)
                style.paragraph_format.space_after = Pt(5 if name != 'Title' else 12)
            style.paragraph_format.first_line_indent = Inches(0)
            style_properties = style.element.get_or_add_pPr()
            inherited_border = style_properties.find(qn('w:pBdr'))
            if inherited_border is not None:
                style_properties.remove(inherited_border)
            if resolved_profile == 'academic':
                r_pr = style.element.get_or_add_rPr()
                for inherited_name in ('spacing', 'kern'):
                    inherited = r_pr.find(qn(f'w:{inherited_name}'))
                    if inherited is not None:
                        r_pr.remove(inherited)

        for list_name in ('List Bullet', 'List Number'):
            list_style = doc.styles[list_name]
            self._set_docx_font(list_style.font, list_style.element, body_font)
            list_style.font.size = Pt(body_size)
            list_style.paragraph_format.left_indent = Inches(0.28 if resolved_profile == 'resume' else 0.32)
            list_style.paragraph_format.first_line_indent = Inches(-0.18)
            list_style.paragraph_format.space_after = Pt(2 if resolved_profile == 'resume' else (0 if resolved_profile == 'academic' else 3))
            if resolved_profile == 'academic':
                list_style.paragraph_format.line_spacing = 2.0

        doc.core_properties.title = resolved_title; doc.core_properties.author = 'Ask Crump'
        doc.core_properties.subject = f'{resolved_profile.title()} document'
        doc.core_properties.keywords = f'{resolved_profile}, professionally formatted, editable document'
        if resolved_profile != 'resume':
            footer = section.footer.paragraphs[0]
            self._add_page_number(footer)
            for run in footer.runs:
                run.font.name = body_font; run.font.size = Pt(8)
                run.font.color.rgb = DocxRGBColor.from_string(MUTED)

        title_p = doc.add_paragraph(style='Title')
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER if resolved_profile == 'academic' else WD_ALIGN_PARAGRAPH.LEFT
        title_p.paragraph_format.first_line_indent = Inches(0); self._add_docx_inline(title_p, resolved_title)

        skipped = False; academic_references = False
        for kind, content in self.blocks(markdown):
            if not skipped and kind in {'h1', 'h2', 'h3', 'h4'} and self._same_title(str(content), resolved_title):
                skipped = True; continue
            if kind in {'h1', 'h2', 'h3', 'h4'}:
                p = doc.add_heading('', level=min(3, max(1, int(kind[1:]))))
                heading_text = str(content).upper() if resolved_profile == 'resume' and kind == 'h1' else str(content)
                self._add_docx_inline(p, heading_text)
                if resolved_profile == 'academic':
                    academic_references = self._clean_inline(str(content)).casefold() in {'references', 'works cited', 'bibliography'}
            elif kind in {'li', 'ol'}:
                p = doc.add_paragraph(style='List Bullet' if kind == 'li' else 'List Number')
                self._add_docx_inline(p, str(content))
            elif kind == 'quote':
                p = doc.add_paragraph(style='Quote'); p.paragraph_format.first_line_indent = Inches(0)
                self._add_docx_inline(p, str(content))
            elif kind == 'code':
                p = doc.add_paragraph(); p.paragraph_format.first_line_indent = Inches(0)
                p.paragraph_format.left_indent = Inches(0.18); p.paragraph_format.space_after = Pt(8)
                run = p.add_run(str(content)); run.font.name = 'Consolas'; run.font.size = Pt(8.7)
                self._shade_paragraph(p, 'F1F3F6')
            elif kind == 'hr':
                p = doc.add_paragraph(); p.paragraph_format.first_line_indent = Inches(0)
                borders = OxmlElement('w:pBdr'); bottom = OxmlElement('w:bottom')
                for key, value in {'val': 'single', 'sz': '6', 'space': '1', 'color': LINE}.items():
                    bottom.set(qn(f'w:{key}'), value)
                borders.append(bottom); p._p.get_or_add_pPr().append(borders)
            elif kind == 'table':
                rows = content
                if not rows: continue
                width = max(len(row) for row in rows); table = doc.add_table(rows=len(rows), cols=width)
                table.autofit = False; table.alignment = WD_TABLE_ALIGNMENT.LEFT
                ratios = self._column_ratios(rows, width)
                usable_inches = 8.5 - margins[1] - margins[3]
                for column_index, ratio in enumerate(ratios):
                    column_width = Inches(usable_inches * ratio)
                    table.columns[column_index].width = column_width
                    for cell in table.columns[column_index].cells:
                        cell.width = column_width
                table.rows[0]._tr.get_or_add_trPr().append(OxmlElement('w:tblHeader'))
                for row_index, values in enumerate(rows):
                    row_properties = table.rows[row_index]._tr.get_or_add_trPr()
                    if row_properties.find(qn('w:cantSplit')) is None:
                        row_properties.append(OxmlElement('w:cantSplit'))
                    for column_index in range(width):
                        cell = table.cell(row_index, column_index); cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
                        self._set_cell_margins(cell)
                        if row_index == 0: self._set_cell_shading(cell, SOFT)
                        elif row_index % 2 == 0: self._set_cell_shading(cell, 'F8F9FA')
                        p = cell.paragraphs[0]; p.paragraph_format.first_line_indent = Inches(0); p.paragraph_format.space_after = Pt(0)
                        p.paragraph_format.keep_together = True
                        p.paragraph_format.keep_with_next = row_index == 0 and len(rows) > 1
                        self._add_docx_inline(p, str(values[column_index] if column_index < len(values) else ''))
                        for run in p.runs:
                            run.font.size = Pt(9.2)
                            if row_index == 0: run.bold = True; run.font.color.rgb = DocxRGBColor.from_string(INK)
                        cell_props = cell._tc.get_or_add_tcPr()
                        cell_borders = OxmlElement('w:tcBorders')
                        for edge in ('top', 'bottom'):
                            border = OxmlElement(f'w:{edge}')
                            border.set(qn('w:val'), 'single')
                            border.set(qn('w:sz'), '4' if row_index == 0 else '2')
                            border.set(qn('w:color'), SLATE if row_index == 0 else LINE)
                            cell_borders.append(border)
                        cell_props.append(cell_borders)
                doc.add_paragraph().paragraph_format.space_after = Pt(0)
            elif content:
                p = doc.add_paragraph()
                if resolved_profile != 'academic':
                    p.paragraph_format.first_line_indent = Inches(0)
                elif academic_references:
                    p.paragraph_format.left_indent = Inches(0.5)
                    p.paragraph_format.first_line_indent = Inches(-0.5)
                self._add_docx_inline(p, str(content))
        out = BytesIO(); doc.save(out); return out.getvalue()

    @staticmethod
    def _xml(text: str) -> str:
        return str(text).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

    @classmethod
    def _rl_markup(cls, text: str) -> str:
        value = cls._xml(text)
        value = re.sub(r'\*\*(.+?)\*\*|__(.+?)__', lambda m: f'<b>{m.group(1) or m.group(2)}</b>', value)
        value = re.sub(r'(?<!\*)\*(.+?)\*(?!\*)|(?<!_)_(.+?)_(?!_)', lambda m: f'<i>{m.group(1) or m.group(2)}</i>', value)
        value = re.sub(r'`(.+?)`', r'<font name="Courier">\1</font>', value)
        value = re.sub(r'\[([^]]+)\]\(([^)]+)\)', r'<u>\1</u>', value)
        return value

    def pdf(self, markdown: str, *, profile: str | None = None, title: str | None = None) -> bytes:
        resolved_title = (title or self.title_from(markdown)).strip()[:160]
        resolved_profile = profile or self.profile_for(markdown); academic = resolved_profile == 'academic'
        out = BytesIO(); margin = 1.0 if academic else 0.76
        document = SimpleDocTemplate(
            out, pagesize=LETTER, leftMargin=margin * inch, rightMargin=margin * inch,
            topMargin=(0.86 if academic else 0.72) * inch, bottomMargin=(0.82 if academic else 0.7) * inch,
            title=resolved_title, author='Ask Crump', subject=f'{resolved_profile.title()} document',
        )
        sample = getSampleStyleSheet(); body_font = 'Times-Roman' if academic else 'Helvetica'
        body = ParagraphStyle(
            'CrumpBody', parent=sample['BodyText'], fontName=body_font, fontSize=12 if academic else 10.5,
            leading=24 if academic else 15, textColor=colors.HexColor(f'#{INK}'), alignment=TA_LEFT,
            spaceAfter=0 if academic else 7, firstLineIndent=0.5 * inch if academic else 0,
        )
        list_body = ParagraphStyle('CrumpListBody', parent=body, firstLineIndent=0, spaceAfter=0 if academic else 2)
        title_style = ParagraphStyle(
            'CrumpTitle', parent=sample['Title'], fontName='Times-Bold' if academic else 'Helvetica-Bold',
            fontSize=12 if academic else 23, leading=24 if academic else 28,
            textColor=colors.HexColor(f'#{INK if academic else SLATE}'),
            alignment=TA_CENTER if academic else TA_LEFT, spaceAfter=0 if academic else 16,
        )
        h1 = ParagraphStyle(
            'CrumpH1', parent=sample['Heading1'], fontName='Times-Bold' if academic else 'Helvetica-Bold',
            fontSize=12 if academic else 15, leading=24 if academic else 19,
            textColor=colors.HexColor(f'#{INK if academic else SLATE}'),
            spaceBefore=0 if academic else 14, spaceAfter=0 if academic else 6, keepWithNext=True,
        )
        h2 = ParagraphStyle(
            'CrumpH2', parent=h1, fontName='Times-Bold' if academic else 'Helvetica-Bold',
            fontSize=12 if academic else 12.5, leading=24 if academic else 16,
        )
        h3 = ParagraphStyle(
            'CrumpH3', parent=h2, fontName='Times-BoldItalic' if academic else 'Helvetica-Bold',
        )
        reference_style = ParagraphStyle(
            'CrumpReference', parent=body, alignment=TA_LEFT, leftIndent=0.5 * inch,
            firstLineIndent=-0.5 * inch,
        )
        code_style = ParagraphStyle('CrumpCode', parent=body, fontName='Courier', fontSize=8.6, leading=11, backColor=colors.HexColor('#F1F3F6'), borderPadding=7, firstLineIndent=0, spaceAfter=8)
        quote_style = ParagraphStyle('CrumpQuote', parent=body, fontName='Times-Italic' if academic else 'Helvetica-Oblique', leftIndent=22, rightIndent=12, firstLineIndent=0, textColor=colors.HexColor(f'#{MUTED}'), borderColor=colors.HexColor(f'#{LINE}'), borderWidth=0.75, borderPadding=7, spaceAfter=10)
        story: list[Any] = [Paragraph(self._rl_markup(resolved_title), title_style)]
        pending: list[ListItem] = []; pending_type = 'bullet'
        def flush_lists() -> None:
            nonlocal pending
            if pending:
                story.append(ListFlowable(
                    pending, bulletType=pending_type, leftIndent=26 if academic else 20,
                    bulletFontName='Times-Roman' if academic else 'Helvetica',
                    bulletFontSize=12 if academic else 7,
                )); story.append(Spacer(1, 5)); pending = []
        skipped = False; academic_references = False
        for kind, content in self.blocks(markdown):
            if not skipped and kind in {'h1', 'h2', 'h3', 'h4'} and self._same_title(str(content), resolved_title):
                skipped = True; continue
            if kind in {'li', 'ol'}:
                desired = '1' if kind == 'ol' else 'bullet'
                if pending and pending_type != desired: flush_lists()
                pending_type = desired; pending.append(ListItem(Paragraph(self._rl_markup(str(content)), list_body), leftIndent=10)); continue
            flush_lists()
            if kind == 'h1':
                story.append(Paragraph(self._rl_markup(str(content)), h1))
            elif kind == 'h2':
                story.append(Paragraph(self._rl_markup(str(content)), h2))
                if academic:
                    academic_references = self._clean_inline(str(content)).casefold() in {'references', 'works cited', 'bibliography'}
            elif kind in {'h3', 'h4'}:
                story.append(Paragraph(self._rl_markup(str(content)), h3))
            elif kind == 'code': story.append(Paragraph(self._xml(str(content)).replace('\n', '<br/>'), code_style))
            elif kind == 'quote': story.append(Paragraph(self._rl_markup(str(content)), quote_style))
            elif kind == 'hr':
                story.extend([Spacer(1, 5), Table([['']], colWidths=[document.width], rowHeights=[0.5], style=TableStyle([('BACKGROUND', (0, 0), (-1, -1), colors.HexColor(f'#{LINE}'))])), Spacer(1, 7)])
            elif kind == 'table' and content:
                width = max(len(row) for row in content); table_rows = []
                for row in content:
                    padded = [*row, *([''] * (width - len(row)))]
                    table_rows.append([Paragraph(self._rl_markup(str(value)), list_body) for value in padded])
                ratios = self._column_ratios(content, width)
                table = Table(table_rows, colWidths=[document.width * ratio for ratio in ratios], repeatRows=1, hAlign='LEFT')
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(f'#{SOFT}')), ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor(f'#{INK}')),
                    ('FONTNAME', (0, 0), (-1, 0), 'Times-Bold' if academic else 'Helvetica-Bold'), ('FONTSIZE', (0, 0), (-1, -1), 8.8),
                    ('LEADING', (0, 0), (-1, -1), 11), ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('LINEBELOW', (0, 0), (-1, 0), 0.8, colors.HexColor(f'#{SLATE}')),
                    ('LINEBELOW', (0, 1), (-1, -1), 0.3, colors.HexColor(f'#{LINE}')),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F8F9FA')]),
                    ('LEFTPADDING', (0, 0), (-1, -1), 6), ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                    ('TOPPADDING', (0, 0), (-1, -1), 6), ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ])); story.extend([table, Spacer(1, 10)])
            elif content: story.append(Paragraph(self._rl_markup(str(content)), reference_style if academic_references else body))
        flush_lists()
        def footer(canvas: Any, doc_template: Any) -> None:
            canvas.saveState(); canvas.setFont('Helvetica', 7.5); canvas.setFillColor(colors.HexColor(f'#{MUTED}'))
            label = str(canvas.getPageNumber())
            canvas.drawRightString(doc_template.pagesize[0] - doc_template.rightMargin, 0.42 * inch, label); canvas.restoreState()
        document.build(story, onFirstPage=footer, onLaterPages=footer); return out.getvalue()

    @staticmethod
    def _ppt_add_text(
        slide: Any,
        text: str,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        size: float,
        color: str,
        bold: bool = False,
        font: str = 'Aptos',
        align: Any = PP_ALIGN.LEFT,
    ) -> Any:
        box = slide.shapes.add_textbox(
            PptInches(left), PptInches(top), PptInches(width), PptInches(height),
        )
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.TOP
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = align
        paragraph.space_after = PptPt(0)
        paragraph.font.name = font
        paragraph.font.size = PptPt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor.from_string(color)
        return box

    @staticmethod
    def _ppt_background(slide: Any, color: str = INK) -> None:
        fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = RGBColor.from_string(color)

    def _ppt_footer(self, slide: Any, number: int, *, dark: bool = False) -> None:
        rule_color = '3A404B' if dark else LINE
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            PptInches(0.72),
            PptInches(7.04),
            PptInches(11.9),
            PptInches(0.012),
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor.from_string(rule_color)
        rule.line.fill.background()
        self._ppt_add_text(
            slide,
            str(number),
            12.06,
            7.12,
            0.55,
            0.16,
            size=8,
            color=SKY if dark else MUTED,
            align=PP_ALIGN.RIGHT,
        )

    @staticmethod
    def _ppt_shape(
        slide: Any,
        shape_type: Any,
        left: float,
        top: float,
        width: float,
        height: float,
        *,
        fill: str,
        line: str | None = None,
    ) -> Any:
        shape = slide.shapes.add_shape(
            shape_type,
            PptInches(left),
            PptInches(top),
            PptInches(width),
            PptInches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill)
        if line:
            shape.line.color.rgb = RGBColor.from_string(line)
        else:
            shape.line.fill.background()
        return shape

    def _ppt_slide_title(
        self,
        slide: Any,
        heading: str,
        *,
        dark: bool = False,
        accent: str = GOLD,
    ) -> None:
        self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 0.72, 0.48, 0.72, 0.055, fill=accent)
        self._ppt_add_text(
            slide,
            'EXECUTIVE BRIEFING',
            0.72,
            0.58,
            2.2,
            0.18,
            size=8.5,
            color=accent,
            bold=True,
        )
        # Keep the hierarchy presentation-scale even when the model returns a
        # long takeaway. Two lines at 35pt are preferable to a document-like
        # title that shrinks into body-copy territory.
        title_size = 38 if len(heading) <= 54 else 35
        self._ppt_add_text(
            slide,
            heading,
            0.72,
            0.86,
            11.7,
            1.08,
            size=title_size,
            color=WHITE if dark else INK,
            bold=True,
            font='Aptos Display',
        )

    @classmethod
    def _ppt_chart_data(
        cls,
        rows: list[list[str]],
    ) -> tuple[list[str], list[tuple[str, list[float], str]]] | None:
        """Return conservative chart data only when a Markdown table is truly numeric."""
        if len(rows) < 4 or len(rows) > 9:
            return None
        width = min(5, max((len(row) for row in rows), default=0))
        if width < 2:
            return None
        padded = [[str(value).strip() for value in [*row, *([''] * (width - len(row)))][:width]] for row in rows]
        categories = [cls._clean_inline(row[0])[:36] for row in padded[1:]]
        if any(not value for value in categories):
            return None

        series: list[tuple[str, list[float], str]] = []
        for column in range(1, width):
            raw_values = [row[column] for row in padded[1:]]
            values: list[float] = []
            formats: list[str] = []
            for raw in raw_values:
                text = cls._clean_inline(raw).strip()
                if not text or text.startswith('='):
                    values = []
                    break
                negative = text.startswith('(') and text.endswith(')')
                percent = text.endswith('%')
                currency = bool(re.match(r'^[\s$£€¥]', text))
                suffix = re.search(r'([kmb])\s*%?$', text, re.I)
                cleaned = re.sub(r'[^0-9.\-]', '', text)
                if not re.fullmatch(r'-?\d+(?:\.\d+)?', cleaned):
                    values = []
                    break
                number = float(cleaned)
                if negative:
                    number = -abs(number)
                if suffix:
                    number *= {'k': 1_000, 'm': 1_000_000, 'b': 1_000_000_000}[suffix.group(1).lower()]
                if percent:
                    number /= 100
                values.append(number)
                formats.append('0%' if percent else ('$0.0,,' if currency and abs(number) >= 1_000_000 else ('$0' if currency else '0.0')))
            if values:
                number_format = formats[0] if len(set(formats)) == 1 else '0.0'
                series.append((cls._clean_inline(padded[0][column])[:42] or f'Series {column}', values, number_format))
        return (categories, series) if series else None

    def _ppt_add_chart(
        self,
        slide: Any,
        rows: list[list[str]],
        *,
        accent: str,
        top: float = 2.05,
        height: float = 4.52,
    ) -> bool:
        prepared = self._ppt_chart_data(rows)
        if not prepared:
            return False
        categories, series_values = prepared
        data = ChartData()
        data.categories = categories
        for label, values, _number_format in series_values:
            data.add_series(label, values)
        use_bar = len(categories) > 5 or max((len(value) for value in categories), default=0) > 14
        chart_type = XL_CHART_TYPE.BAR_CLUSTERED if use_bar else XL_CHART_TYPE.COLUMN_CLUSTERED
        chart = slide.shapes.add_chart(
            chart_type,
            PptInches(0.82),
            PptInches(top),
            PptInches(11.7),
            PptInches(height),
            data,
        ).chart
        chart.has_title = False
        chart.has_legend = len(series_values) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.RIGHT if use_bar else XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.name = 'Aptos'
            chart.legend.font.size = PptPt(10)
            chart.legend.font.color.rgb = RGBColor.from_string(MUTED)
        palette = [accent, TEAL, COPPER, SKY]
        for index, series in enumerate(chart.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(palette[index % len(palette)])
            series.format.line.fill.background()
        chart.category_axis.tick_labels.font.name = 'Aptos'
        chart.category_axis.tick_labels.font.size = PptPt(10.5)
        chart.category_axis.tick_labels.font.color.rgb = RGBColor.from_string(INK)
        chart.value_axis.tick_labels.font.name = 'Aptos'
        chart.value_axis.tick_labels.font.size = PptPt(9.5)
        chart.value_axis.tick_labels.font.color.rgb = RGBColor.from_string(MUTED)
        formats = {number_format for _label, _values, number_format in series_values}
        if len(formats) == 1:
            chart.value_axis.tick_labels.number_format = next(iter(formats))
            chart.value_axis.tick_labels.number_format_is_linked = False
        chart.value_axis.has_major_gridlines = True
        return True

    @staticmethod
    def _ppt_normalize_axis_ids(data: bytes) -> bytes:
        """Keep native charts readable by strict OOXML consumers as well as PowerPoint."""
        source = BytesIO(data)
        output = BytesIO()
        with ZipFile(source, 'r') as incoming, ZipFile(output, 'w', ZIP_DEFLATED) as outgoing:
            for info in incoming.infolist():
                payload = incoming.read(info.filename)
                if info.filename.startswith('ppt/charts/chart') and info.filename.endswith('.xml'):
                    xml = payload.decode('utf-8')
                    xml = re.sub(
                        r'(<c:(?:axId|crossAx)\s+val=")-(\d+)(")',
                        lambda match: f'{match.group(1)}{match.group(2)}{match.group(3)}',
                        xml,
                    )
                    payload = xml.encode('utf-8')
                outgoing.writestr(info, payload)
        return output.getvalue()

    def pptx(self, markdown: str, *, title: str | None = None) -> bytes:
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)
        resolved_title = (title or self.title_from(markdown)).strip()[:160]
        prs.core_properties.title = resolved_title
        prs.core_properties.author = 'Ask Crump'
        cover = prs.slides.add_slide(prs.slide_layouts[6])
        self._ppt_background(cover, INK)
        self._ppt_shape(cover, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, 7.5, fill=TEAL)
        self._ppt_shape(cover, MSO_SHAPE.RECTANGLE, 0.82, 0.84, 0.78, 0.06, fill=GOLD)
        ring = cover.shapes.add_shape(
            MSO_SHAPE.OVAL,
            PptInches(10.72),
            PptInches(-0.62),
            PptInches(3.35),
            PptInches(3.35),
        )
        ring.fill.background()
        ring.line.color.rgb = RGBColor.from_string(GOLD)
        ring.line.width = PptPt(2)
        self._ppt_shape(cover, MSO_SHAPE.OVAL, 11.7, 0.37, 0.34, 0.34, fill=COPPER)
        self._ppt_add_text(
            cover,
            'EXECUTIVE PRESENTATION',
            0.82,
            1.02,
            3.2,
            0.22,
            size=9,
            color=SKY,
            bold=True,
        )
        title_size = 54 if len(resolved_title) <= 60 else (45 if len(resolved_title) <= 95 else 38)
        self._ppt_add_text(
            cover,
            resolved_title,
            0.82,
            1.45,
            10.6,
            2.3,
            size=title_size,
            color=WHITE,
            bold=True,
            font='Aptos Display',
        )
        subtitle = 'Clear thinking. Structured for decisions.'
        for kind, content in self.blocks(markdown):
            if kind == 'p' and self._clean_inline(str(content)) != self._clean_inline(resolved_title):
                subtitle = self._clean_inline(str(content))[:170]
                break
        self._ppt_add_text(cover, subtitle, 0.86, 4.55, 9.55, 0.85, size=19, color=CREAM)
        self._ppt_shape(cover, MSO_SHAPE.RECTANGLE, 0.84, 6.52, 3.25, 0.045, fill=TEAL)
        self._ppt_shape(cover, MSO_SHAPE.RECTANGLE, 4.09, 6.52, 1.15, 0.045, fill=GOLD)
        self._ppt_shape(cover, MSO_SHAPE.RECTANGLE, 5.24, 6.52, 0.52, 0.045, fill=COPPER)

        groups: list[tuple[str, list[Any]]] = []; current_title = 'Executive overview'; current: list[Any] = []; first_heading = False
        for kind, content in self.blocks(markdown):
            if kind in {'h1', 'h2', 'h3', 'h4'}:
                clean = self._clean_inline(str(content))
                if not first_heading and self._same_title(clean, resolved_title):
                    first_heading = True; continue
                if current: groups.append((current_title, current))
                current_title, current, first_heading = clean, [], True
            elif kind != 'hr': current.append((kind, content))
        if current: groups.append((current_title, current))
        if not groups: groups = [('Executive overview', [('p', 'A clear, decision-ready summary.')])]

        content_number = 0
        accents = [TEAL, GOLD, COPPER, SKY]
        limited_groups = groups[:18]
        for group_index, (heading, items) in enumerate(limited_groups):
            text_items: list[str] = []
            table_items: list[list[list[str]]] = []
            for kind, content in items:
                if kind == 'table':
                    table_items.append(content)
                else:
                    text_items.append(self._clean_inline(str(content).replace('\n', ' ')))
            # A single sentence immediately before a table is its explanatory
            # lead, not a reason to create a sparse duplicate-topic slide.
            table_lead = text_items[0] if table_items and len(text_items) == 1 and text_items[0] else None
            paged_text_items = [] if table_lead else text_items
            for start in range(0, len(paged_text_items), 5):
                page_items = [item for item in paged_text_items[start:start + 5] if item]
                if not page_items:
                    continue
                content_number += 1
                accent = accents[(content_number - 1) % len(accents)]
                is_closing = group_index == len(limited_groups) - 1 and not table_items
                dark = is_closing or content_number % 4 == 0
                background = NAVY if dark else PAPER
                body_color = CREAM if dark else INK
                secondary_color = SKY if dark else MUTED
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._ppt_background(slide, background)
                display = heading if start == 0 else f'{heading} — continued'
                self._ppt_slide_title(slide, display, dark=dark, accent=accent)
                if len(page_items) == 1:
                    watermark = '303746' if dark else SOFT
                    self._ppt_add_text(
                        slide,
                        f'{content_number:02d}',
                        9.1,
                        2.0,
                        2.45,
                        1.35,
                        size=82,
                        color=watermark,
                        bold=True,
                        font='Aptos Display',
                        align=PP_ALIGN.RIGHT,
                    )
                    self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 2.12, 0.07, 3.05, fill=accent)
                    self._ppt_add_text(
                        slide,
                        f'{content_number:02d}',
                        1.2,
                        2.06,
                        0.82,
                        0.5,
                        size=13,
                        color=accent,
                        bold=True,
                    )
                    size = 31 if len(page_items[0]) <= 240 else 25
                    self._ppt_add_text(
                        slide,
                        page_items[0][:560],
                        1.2,
                        2.63,
                        7.55,
                        2.8,
                        size=size,
                        color=body_color,
                        font='Aptos Display',
                    )
                elif len(page_items) == 2 and not dark:
                    dark_on_right = content_number % 2 == 1
                    panel_x = 7.08 if dark_on_right else 0
                    self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, panel_x, 1.92, 6.253, 5.08, fill=NAVY)
                    for item_index, item in enumerate(page_items):
                        x = 0.82 if item_index == 0 else (7.7 if dark_on_right else 7.02)
                        on_dark_panel = item_index == (1 if dark_on_right else 0)
                        number_color = GOLD if on_dark_panel else accent
                        text_color = CREAM if on_dark_panel else INK
                        self._ppt_add_text(
                            slide,
                            f'{item_index + 1:02d}',
                            x,
                            2.26,
                            0.7,
                            0.34,
                            size=13,
                            color=number_color,
                            bold=True,
                        )
                        self._ppt_shape(
                            slide,
                            MSO_SHAPE.RECTANGLE,
                            x,
                            2.72,
                            0.72,
                            0.055,
                            fill=number_color,
                        )
                        self._ppt_add_text(
                            slide,
                            item[:430],
                            x,
                            3.12,
                            5.1 if item_index == 0 else 4.82,
                            2.55,
                            size=23 if len(item) <= 230 else 18,
                            color=text_color,
                            font='Aptos Display',
                        )
                elif len(page_items) == 3 and not dark:
                    self._ppt_add_text(
                        slide,
                        '01',
                        0.82,
                        2.14,
                        0.72,
                        0.35,
                        size=13,
                        color=accent,
                        bold=True,
                    )
                    self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 2.67, 0.72, 0.055, fill=accent)
                    self._ppt_add_text(
                        slide,
                        page_items[0][:430],
                        0.82,
                        3.04,
                        4.75,
                        2.65,
                        size=25 if len(page_items[0]) <= 230 else 20,
                        color=body_color,
                        font='Aptos Display',
                    )
                    self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 5.86, 2.02, 0.012, 3.98, fill=LINE)
                    for item_index, item in enumerate(page_items[1:], start=2):
                        y = 2.16 + (item_index - 2) * 1.84
                        self._ppt_add_text(
                            slide,
                            f'{item_index:02d}',
                            6.28,
                            y,
                            0.62,
                            0.28,
                            size=11,
                            color=accent,
                            bold=True,
                        )
                        self._ppt_add_text(
                            slide,
                            item[:360],
                            7.15,
                            y - 0.02,
                            5.08,
                            1.3,
                            size=19 if len(item) <= 230 else 17,
                            color=body_color,
                        )
                        if item_index == 2:
                            self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 6.28, 3.65, 5.85, 0.012, fill=LINE)
                elif len(page_items) <= 3:
                    count = len(page_items)
                    gap = 0.34
                    column_width = (11.9 - gap * (count - 1)) / count
                    for item_index, item in enumerate(page_items):
                        x = 0.72 + item_index * (column_width + gap)
                        if item_index:
                            self._ppt_shape(
                                slide,
                                MSO_SHAPE.RECTANGLE,
                                x - gap / 2,
                                2.0,
                                0.012,
                                3.95,
                                fill='3A404B' if dark else LINE,
                            )
                        self._ppt_shape(slide, MSO_SHAPE.OVAL, x, 2.08, 0.42, 0.42, fill=accent)
                        self._ppt_add_text(
                            slide,
                            f'{item_index + 1:02d}',
                            x,
                            2.17,
                            0.42,
                            0.18,
                            size=8,
                            color=WHITE,
                            bold=True,
                            align=PP_ALIGN.CENTER,
                        )
                        size = 20 if len(item) <= 230 else 17
                        self._ppt_add_text(
                            slide,
                            item[:390],
                            x,
                            2.78,
                            column_width - 0.12,
                            2.95,
                            size=size,
                            color=body_color,
                        )
                else:
                    for item_index, item in enumerate(page_items):
                        column = item_index % 2
                        row = item_index // 2
                        x = 0.76 + column * 6.08
                        y = 2.03 + row * 1.38
                        self._ppt_shape(
                            slide,
                            MSO_SHAPE.RECTANGLE,
                            x,
                            y + 0.03,
                            0.07,
                            0.74,
                            fill=accent if row % 2 == 0 else secondary_color,
                        )
                        self._ppt_add_text(
                            slide,
                            f'{item_index + 1:02d}',
                            x + 0.25,
                            y + 0.03,
                            0.42,
                            0.25,
                            size=9,
                            color=accent,
                            bold=True,
                        )
                        size = 17 if len(item) <= 230 else 16
                        self._ppt_add_text(
                            slide,
                            item[:300],
                            x + 0.82,
                            y - 0.02,
                            4.92,
                            1.04,
                            size=size,
                            color=body_color,
                        )
                self._ppt_footer(slide, len(prs.slides), dark=dark)
            for table_index, table_rows in enumerate(table_items):
                if not table_rows:
                    continue
                chunks = [table_rows[:8]]
                if len(table_rows) > 8:
                    chunks.extend([[table_rows[0], *table_rows[index:index + 7]] for index in range(8, len(table_rows), 7)])
                for chunk_index, rows in enumerate(chunks[:3]):
                    content_number += 1
                    accent = accents[(content_number - 1) % len(accents)]
                    width = min(6, max(len(row) for row in rows))
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    self._ppt_background(slide, PAPER)
                    display = heading if chunk_index == 0 else f'{heading} — continued'
                    lead = table_lead if table_index == 0 and chunk_index == 0 else None
                    chart_top = 2.48 if lead else 2.05
                    chart_height = 4.08 if lead else 4.52
                    if chunk_index == 0 and self._ppt_add_chart(
                        slide,
                        rows,
                        accent=accent,
                        top=chart_top,
                        height=chart_height,
                    ):
                        # Charts are native and remain editable. Repaint the title band after the
                        # chart so strict renderers cannot obscure the slide headline with the
                        # chart canvas.
                        title_band_height = 2.38 if lead else 1.82
                        self._ppt_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, title_band_height, fill=PAPER)
                        self._ppt_slide_title(slide, display, accent=accent)
                        if lead:
                            self._ppt_add_text(
                                slide,
                                lead[:260],
                                0.82,
                                1.78,
                                9.35,
                                0.42,
                                size=16,
                                color=MUTED,
                            )
                        self._ppt_add_text(
                            slide,
                            'DATA VIEW',
                            10.15,
                            1.84 if lead else 1.52,
                            2.35,
                            0.2,
                            size=8,
                            color=accent,
                            bold=True,
                            align=PP_ALIGN.RIGHT,
                        )
                        self._ppt_footer(slide, len(prs.slides))
                        continue
                    self._ppt_slide_title(slide, display, accent=accent)
                    table_top = 2.48 if lead else 1.92
                    if lead:
                        self._ppt_add_text(
                            slide,
                            lead[:260],
                            0.82,
                            1.78,
                            11.1,
                            0.42,
                            size=16,
                            color=MUTED,
                        )
                    table_height = max(1.6, min(4.2 if lead else 4.75, 0.58 * len(rows) + 0.3))
                    table = slide.shapes.add_table(
                        len(rows),
                        width,
                        PptInches(0.72),
                        PptInches(table_top),
                        PptInches(11.9),
                        PptInches(table_height),
                    ).table
                    ratios = self._column_ratios(rows, width)
                    for col_index, ratio in enumerate(ratios):
                        table.columns[col_index].width = PptInches(11.9 * ratio)
                    for row_index, row in enumerate(rows):
                        for col_index in range(width):
                            cell = table.cell(row_index, col_index)
                            cell.text = self._clean_inline(str(row[col_index] if col_index < len(row) else ''))[:160]
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor.from_string(
                                SLATE if row_index == 0 else (MIST if row_index % 2 == 0 else WHITE),
                            )
                            cell.margin_left = cell.margin_right = PptInches(0.1)
                            cell.margin_top = cell.margin_bottom = PptInches(0.06)
                            p = cell.text_frame.paragraphs[0]
                            p.font.name = 'Aptos'
                            p.font.size = PptPt(13.5)
                            p.font.bold = row_index == 0 or col_index == 0
                            p.font.color.rgb = RGBColor.from_string(
                                WHITE if row_index == 0 else (TEAL if col_index == 0 else INK),
                            )
                    self._ppt_footer(slide, len(prs.slides))
        out = BytesIO()
        prs.save(out)
        return self._ppt_normalize_axis_ids(out.getvalue())

    @staticmethod
    def _safe_sheet_name(value: str, fallback: str) -> str:
        name = re.sub(r'[\\/*?:\[\]]+', ' ', str(value or '')).strip(); name = re.sub(r'\s+', ' ', name)[:31]
        return name or fallback

    @staticmethod
    def _safe_table_name(value: str, index: int) -> str:
        stem = re.sub(r'[^A-Za-z0-9_]+', '_', str(value or '')).strip('_')
        if not stem or stem[0].isdigit(): stem = f'Table_{stem}'
        return f'{stem[:200]}_{index}'

    @staticmethod
    def _safe_formula(value: str) -> bool:
        formula = str(value or '').strip()
        if not formula.startswith('=') or len(formula) > 500: return False
        if re.search(r"\b(HYPERLINK|WEBSERVICE|FILTERXML|RTD|DDE|CALL|REGISTER)\s*\(", formula, re.I): return False
        if re.search(r"https?://|file:|cmd\||powershell|\[[^]]+\]|'[^']+'!", formula, re.I): return False
        return bool(re.fullmatch(r"=[A-Za-z0-9_.$!(),:+\-*/^%<>=\s\"]+", formula))

    @classmethod
    def _typed_cell(cls, value: Any) -> tuple[Any, str | None]:
        raw = str(value or '').strip()
        if not raw: return '', None
        if raw.startswith('='): return (raw, None) if cls._safe_formula(raw) else (f"'{raw}", None)
        if re.fullmatch(r'\(?\$\s*[-+]?\d[\d,]*(?:\.\d+)?\)?', raw):
            negative = raw.startswith('(') and raw.endswith(')'); number = float(re.sub(r'[$,()\s]', '', raw))
            return (-number if negative else number), '$#,##0.00;[Red]-$#,##0.00'
        if re.fullmatch(r'[-+]?\d[\d,]*(?:\.\d+)?%', raw): return float(raw.replace(',', '').replace('%', '')) / 100, '0.0%'
        if re.fullmatch(r'[-+]?\d[\d,]*', raw): return int(raw.replace(',', '')), '#,##0'
        if re.fullmatch(r'[-+]?\d[\d,]*\.\d+', raw): return float(raw.replace(',', '')), '#,##0.00'
        for date_format in ('%Y-%m-%d', '%m/%d/%Y', '%m/%d/%y'):
            try: return datetime.strptime(raw, date_format), 'mmm d, yyyy'
            except ValueError: continue
        return cls._clean_inline(raw), None

    @classmethod
    def _tables_with_titles(cls, markdown: str) -> list[tuple[str, list[list[str]]]]:
        output: list[tuple[str, list[list[str]]]] = []; heading = 'Data'
        for kind, content in cls.blocks(markdown):
            if kind in {'h1', 'h2', 'h3', 'h4'}: heading = cls._clean_inline(str(content))
            elif kind == 'table': output.append((heading, content))
        return output

    @classmethod
    def _markdown_tables(cls, lines: list[str]) -> list[list[list[str]]]:
        return [table for _, table in cls._tables_with_titles('\n'.join(lines))]

    def _style_xlsx_sheet(self, sheet: Any) -> None:
        sheet.sheet_view.showGridLines = False; sheet.freeze_panes = 'A2'; sheet.print_title_rows = '1:1'
        sheet.sheet_properties.pageSetUpPr.fitToPage = True; sheet.page_setup.orientation = 'landscape'
        sheet.page_setup.fitToWidth = 1; sheet.page_setup.fitToHeight = 0
        sheet.print_area = sheet.dimensions
        sheet.page_margins.left = sheet.page_margins.right = 0.35
        sheet.page_margins.top = sheet.page_margins.bottom = 0.5
        sheet.oddFooter.center.text = '&[Tab]  •  &[Page] of &[Pages]'; sheet.oddFooter.center.size = 8; sheet.oddFooter.center.color = MUTED
        thin = Side(style='hair', color=LINE)
        for cell in sheet[1]:
            cell.fill = PatternFill('solid', fgColor=SLATE); cell.font = Font(name='Aptos', size=9.5, bold=True, color='FFFFFF')
            cell.alignment = Alignment(vertical='center', wrap_text=True); cell.border = Border(bottom=Side(style='medium', color=SLATE))
        sheet.row_dimensions[1].height = 25
        for row in sheet.iter_rows(min_row=2):
            for cell in row:
                cell.font = Font(name='Aptos', size=9.5, color=INK); cell.alignment = Alignment(vertical='top', wrap_text=True)
                cell.border = Border(bottom=thin)
                if str(cell.value or '').startswith('='):
                    cell.fill = PatternFill('solid', fgColor='F1F5F9')
                    cell.font = Font(name='Aptos', size=9.5, color='274A67')
        for column in sheet.columns:
            letter = column[0].column_letter; values = [len(str(cell.value or '')) for cell in column[:100]]
            sheet.column_dimensions[letter].width = min(40, max(10, max(values, default=8) + 2))
        sheet.sheet_properties.tabColor = SLATE

    @staticmethod
    def _header_number_format(header: str) -> str | None:
        label = str(header or '').strip().casefold()
        if re.search(r'\b(mrr|arr|revenue|price|cost|budget|spend|income|amount|cash|sales|profit)\b', label):
            return '$#,##0.00;[Red]-$#,##0.00'
        if re.search(r'\b(rate|growth|conversion|churn|margin|percentage|percent)\b|%', label):
            return '0.0%'
        if re.search(r'\b(users?|customers?|leads?|count|quantity|units?|months?)\b', label):
            return '#,##0'
        return None

    def xlsx(self, markdown: str, *, title: str | None = None) -> bytes:
        resolved_title = (title or self.title_from(markdown)).strip()[:160]
        workbook = Workbook(); overview = workbook.active; overview.title = 'Overview'; overview.sheet_view.showGridLines = False
        overview.sheet_properties.tabColor = SLATE; overview.merge_cells('A1:C1'); overview['A1'] = resolved_title
        overview['A1'].font = Font(name='Aptos', size=20, bold=True, color=INK)
        overview['A1'].alignment = Alignment(vertical='center'); overview['A1'].border = Border(bottom=Side(style='medium', color=SLATE))
        overview.row_dimensions[1].height = 36
        overview['A3'] = 'WORKBOOK INDEX'; overview['A3'].font = Font(name='Aptos', size=8, bold=True, color=SLATE)
        overview.merge_cells('A5:C5')
        overview['A5'] = 'Use this index to move between the model’s structured tables. Inputs, formulas, units, and number formats remain editable.'
        overview['A5'].alignment = Alignment(wrap_text=True, vertical='top'); overview['A5'].font = Font(name='Aptos', size=10, color=MUTED)
        overview.row_dimensions[5].height = 32
        for column, label in enumerate(('Sheet', 'Purpose', 'Records'), start=1):
            cell = overview.cell(7, column, label); cell.fill = PatternFill('solid', fgColor=SOFT)
            cell.font = Font(name='Aptos', size=9.5, bold=True, color=INK); cell.alignment = Alignment(vertical='center')
            cell.border = Border(bottom=Side(style='thin', color=SLATE))
        overview.row_dimensions[7].height = 22
        overview.column_dimensions['A'].width = 30; overview.column_dimensions['B'].width = 48; overview.column_dimensions['C'].width = 14
        overview.sheet_properties.pageSetUpPr.fitToPage = True; overview.page_setup.orientation = 'portrait'
        overview.page_setup.fitToWidth = 1; overview.page_setup.fitToHeight = 1
        overview.page_margins.left = overview.page_margins.right = 0.55
        overview.page_margins.top = overview.page_margins.bottom = 0.55
        overview.oddFooter.center.text = '&[Page] of &[Pages]'; overview.oddFooter.center.size = 8; overview.oddFooter.center.color = MUTED

        tables = self._tables_with_titles(markdown)
        if tables:
            used_names = {'overview'}
            for index, (heading, rows) in enumerate(tables, start=1):
                base = self._safe_sheet_name(heading, f'Data {index}'); name = base; suffix = 2
                while name.casefold() in used_names:
                    tail = f' {suffix}'; name = f'{base[:31 - len(tail)]}{tail}'; suffix += 1
                used_names.add(name.casefold()); sheet = workbook.create_sheet(name); width = max(len(row) for row in rows)
                headers = [self._clean_inline(str(rows[0][column] if column < len(rows[0]) else f'Column {column + 1}')) or f'Column {column + 1}' for column in range(width)]
                seen: dict[str, int] = {}
                for column, header in enumerate(headers):
                    key = header.casefold(); seen[key] = seen.get(key, 0) + 1
                    if seen[key] > 1: headers[column] = f'{header} {seen[key]}'
                sheet.append(headers)
                for row in rows[1:]:
                    typed = []; formats: dict[int, str] = {}
                    for column in range(width):
                        value, number_format = self._typed_cell(row[column] if column < len(row) else '')
                        typed.append(value)
                        if number_format: formats[column + 1] = number_format
                    sheet.append(typed)
                    for column, number_format in formats.items():
                        sheet.cell(sheet.max_row, column).number_format = number_format
                self._style_xlsx_sheet(sheet)
                for column, header in enumerate(headers, start=1):
                    inferred_format = self._header_number_format(header)
                    if not inferred_format:
                        continue
                    for row_index in range(2, sheet.max_row + 1):
                        cell = sheet.cell(row_index, column)
                        if cell.number_format == 'General' and (
                            isinstance(cell.value, (int, float)) or str(cell.value or '').startswith('=')
                        ):
                            cell.number_format = inferred_format
                if sheet.max_row >= 2:
                    table = XlsxTable(displayName=self._safe_table_name(name, index), ref=sheet.dimensions)
                    table.tableStyleInfo = TableStyleInfo(name='TableStyleMedium2', showFirstColumn=False, showLastColumn=False, showRowStripes=True, showColumnStripes=False)
                    sheet.add_table(table)
                    for column in range(1, sheet.max_column + 1):
                        numeric = [sheet.cell(row, column).value for row in range(2, sheet.max_row + 1)]
                        numeric_values = [value for value in numeric if isinstance(value, (int, float))]
                        if len(numeric_values) >= 6 and max(numeric_values) != min(numeric_values):
                            letter = sheet.cell(1, column).column_letter
                            sheet.conditional_formatting.add(f'{letter}2:{letter}{sheet.max_row}', DataBarRule(start_type='min', end_type='max', color='7B91A6', showValue=True))
                index_row = 7 + index
                overview.cell(index_row, 1, name); overview.cell(index_row, 2, heading); overview.cell(index_row, 3, max(0, len(rows) - 1))
                overview.cell(index_row, 1).hyperlink = f"#'{name}'!A1"; overview.cell(index_row, 1).style = 'Hyperlink'
                for cell in overview[index_row]:
                    cell.font = Font(name='Aptos', size=9.5, color='245B82' if cell.column == 1 else INK, underline='single' if cell.column == 1 else None)
                    cell.alignment = Alignment(vertical='top', wrap_text=True); cell.border = Border(bottom=Side(style='hair', color=LINE))
                overview.row_dimensions[index_row].height = 24
        else:
            sheet = workbook.create_sheet('Brief'); sheet.sheet_view.showGridLines = False
            sheet.column_dimensions['A'].width = 24; sheet.column_dimensions['B'].width = 82; row = 1
            for kind, content in self.blocks(markdown):
                if kind in {'h1', 'h2', 'h3', 'h4'}:
                    sheet.merge_cells(start_row=row, start_column=1, end_row=row, end_column=2); cell = sheet.cell(row, 1, self._clean_inline(str(content)))
                    cell.fill = PatternFill('solid', fgColor=SLATE if row == 1 else 'EDEFF3'); cell.font = Font(name='Aptos', size=16 if row == 1 else 12, bold=True, color='FFFFFF' if row == 1 else INK)
                    cell.alignment = Alignment(vertical='center', wrap_text=True); sheet.row_dimensions[row].height = 30
                elif kind == 'hr':
                    continue
                elif kind != 'table':
                    sheet.cell(row, 1, kind.upper()); sheet.cell(row, 1).font = Font(name='Aptos', size=8, bold=True, color=SLATE)
                    sheet.cell(row, 2, self._clean_inline(str(content))); sheet.cell(row, 2).font = Font(name='Aptos', size=10, color=INK)
                    sheet.cell(row, 2).alignment = Alignment(wrap_text=True, vertical='top'); sheet.row_dimensions[row].height = max(24, min(90, 15 * (1 + len(str(content)) // 90)))
                row += 1
            sheet.freeze_panes = 'A2'; sheet.sheet_properties.tabColor = SLATE
            overview.cell(8, 1, 'Brief'); overview.cell(8, 2, 'Formatted document content'); overview.cell(8, 3, max(0, row - 1))
            overview.cell(8, 1).hyperlink = "#'Brief'!A1"; overview.cell(8, 1).style = 'Hyperlink'
        overview.print_area = f'A1:C{max(10, overview.max_row + 1)}'
        overview.freeze_panes = 'A8'
        workbook.calculation.calcMode = 'auto'; workbook.calculation.fullCalcOnLoad = True; workbook.calculation.forceFullCalc = True
        workbook.properties.title = resolved_title; workbook.properties.creator = 'Ask Crump'; workbook.properties.subject = 'Professionally formatted editable workbook'
        out = BytesIO(); workbook.save(out); return out.getvalue()

    async def create(self, *, user_id: str, markdown: str, format_name: str, chat_id: str | None, message_id: str | None, title: str | None = None, brief: str | None = None, purpose: Any = None) -> dict[str, Any]:
        fmt = self.normalize_format(format_name)
        if not fmt: raise ValueError('Unsupported document format.')
        resolved_title = (title or self.title_from(markdown)).strip()[:160]; profile = self.profile_for(markdown, brief or '', fmt, purpose)
        stem = self.safe_stem(resolved_title)
        if fmt == 'docx': data = self.docx(markdown, profile=profile, title=resolved_title)
        elif fmt == 'pdf': data = self.pdf(markdown, profile=profile, title=resolved_title)
        elif fmt == 'pptx': data = self.pptx(markdown, title=resolved_title)
        elif fmt == 'xlsx': data = self.xlsx(markdown, title=resolved_title)
        elif fmt == 'md': data = markdown.encode('utf-8')
        else: data = self._clean_inline(markdown).encode('utf-8')
        row = await self.files.store_bytes(
            user_id=user_id, data=data, filename=f'{stem}.{fmt}', mime_type=MIME[fmt],
            kind='generated_document', chat_id=chat_id, message_id=message_id,
            metadata={'format': fmt, 'title': resolved_title, 'profile': profile},
        )
        result = self.files.public_file(row); result.update({'format': fmt, 'title': resolved_title, 'profile': profile}); return result
