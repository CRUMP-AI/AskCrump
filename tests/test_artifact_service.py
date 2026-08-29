import asyncio
from io import BytesIO
import zipfile

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from backend.artifact_service import ArtifactService


class DummyFiles:
    def __init__(self):
        self.stored: dict | None = None

    async def store_bytes(self, **kwargs):
        self.stored = kwargs
        return {'id': 'fixture-file', **kwargs}

    @staticmethod
    def public_file(row):
        return {'id': row['id'], 'name': row['filename'], 'size': len(row['data'])}


def service():
    return ArtifactService(DummyFiles())


SAMPLE = """# Product Brief

## Purpose
A concise, polished artifact.

- First point
- Second point

| Item | Value |
| --- | --- |
| Quality | High |
| Readiness | 95% |
| Budget | $12,500 |
| Launch | 2026-09-15 |
"""

NUMERIC_DECK = """# Growth Brief

## Activation is improving across channels

| Channel | Activation | Retention |
| --- | --- | --- |
| Organic | 42% | 31% |
| Referral | 58% | 46% |
| Paid | 36% | 24% |
| Partner | 49% | 39% |
"""

RHYTHM_DECK = """# Operating Brief

## Context
- The opportunity is specific and measurable.
- The current path creates avoidable friction.

## Decision
- Concentrate the experience around one valuable outcome.
- Make the next action obvious.

## Execution
- Assign an owner and a release gate.
- Measure activation and seven-day continuation.
"""


def test_binary_artifacts_are_valid_containers():
    generator = service()
    docx = generator.docx(SAMPLE)
    pptx = generator.pptx(SAMPLE)
    xlsx = generator.xlsx(SAMPLE)
    pdf = generator.pdf(SAMPLE)
    assert zipfile.is_zipfile(BytesIO(docx))
    assert zipfile.is_zipfile(BytesIO(pptx))
    assert zipfile.is_zipfile(BytesIO(xlsx))
    assert pdf.startswith(b'%PDF')


def test_horizontal_rules_render_without_becoming_headings():
    markdown = """# Sample Resume

---

## Experience

- Built a reliable document pipeline.

---

## Education

A finished section.
"""
    generator = service()
    docx = Document(BytesIO(generator.docx(markdown, profile='resume')))
    deck = Presentation(BytesIO(generator.pptx(markdown)))
    workbook = load_workbook(BytesIO(generator.xlsx(markdown)), data_only=False)
    pdf = PdfReader(BytesIO(generator.pdf(markdown)))

    assert any('Experience' in paragraph.text for paragraph in docx.paragraphs)
    assert all(paragraph.text != 'HR' for paragraph in docx.paragraphs)
    visible_slide_text = '\n'.join(
        shape.text for slide in deck.slides for shape in slide.shapes
        if getattr(shape, 'has_text_frame', False)
    )
    assert 'Experience' in visible_slide_text
    assert '\nHR\n' not in f'\n{visible_slide_text}\n'
    sheet_values = [
        str(cell.value or '')
        for sheet in workbook.worksheets
        for row in sheet.iter_rows()
        for cell in row
    ]
    assert 'HR' not in sheet_values
    assert any('Experience' in (page.extract_text() or '') for page in pdf.pages)


def test_supported_format_aliases():
    assert ArtifactService.normalize_format('Word') == 'docx'
    assert ArtifactService.normalize_format('.PDF') == 'pdf'
    assert ArtifactService.normalize_format('powerpoint') == 'pptx'
    assert ArtifactService.normalize_format('excel') == 'xlsx'


def test_natural_language_document_requests_are_detected():
    assert ArtifactService.detect_request(
        'Write the full manuscript and send it as a Word documented file.'
    ) == 'docx'
    assert ArtifactService.detect_request('Please compose the proposal and deliver a PDF.') == 'pdf'
    assert ArtifactService.detect_request('Create a downloadable resume document.') == 'docx'
    assert ArtifactService.detect_request('Could you deliver this in a document format?') == 'docx'
    assert ArtifactService.detect_request('Write a report here in the chat.') is None


def test_book_scale_work_is_handed_to_a_persistent_manuscript():
    assert ArtifactService.is_long_form_request(
        'Write a complete 70,000 word novel and send it as a Word document.'
    ) is True
    assert ArtifactService.is_long_form_request('Draft my full-length memoir from start to finish.') is True
    assert ArtifactService.is_long_form_request('Create a book summary PDF.') is False
    assert ArtifactService.is_long_form_request('Recommend a book about Georgia history.') is False


def test_word_export_has_real_hierarchy_table_metadata_and_neutral_page_number():
    data = service().docx(SAMPLE, profile='business')
    document = Document(BytesIO(data))
    assert document.core_properties.author == 'Ask Crump'
    assert document.core_properties.title == 'Product Brief'
    assert document.sections[0].top_margin.inches < 0.9
    assert document.styles['Title'].font.size.pt >= 24
    assert document.tables
    assert document.tables[0].cell(0, 0).text == 'Item'
    assert document.tables[0].autofit is False
    footer_xml = document.sections[0].footer._element.xml
    assert 'PAGE' in footer_xml
    assert document.sections[0].header.paragraphs[0].text == ''
    assert 'Ask Crump' not in document.sections[0].footer.paragraphs[0].text


def test_word_tables_keep_headers_with_data_and_do_not_split_rows():
    document = Document(BytesIO(service().docx(SAMPLE, profile='business')))
    table = document.tables[0]

    assert len(table.rows) > 1
    assert all('w:cantSplit' in row._tr.xml for row in table.rows)
    assert all(
        paragraph.paragraph_format.keep_with_next
        for cell in table.rows[0].cells
        for paragraph in cell.paragraphs
    )
    assert all(
        not paragraph.paragraph_format.keep_with_next
        for cell in table.rows[1].cells
        for paragraph in cell.paragraphs
    )


def test_academic_and_resume_profiles_use_professional_conventions():
    generator = service()
    assert generator.profile_for('', 'Write a college research essay in APA style', 'docx') == 'academic'
    assert generator.profile_for('', 'Create my résumé for a product role', 'docx') == 'resume'

    academic = Document(BytesIO(generator.docx('# Evidence and Judgment\n\nA finished paragraph.', profile='academic')))
    assert academic.styles['Normal'].font.name == 'Times New Roman'
    assert academic.styles['Normal'].font.size.pt == 12
    assert academic.styles['Normal'].paragraph_format.line_spacing == 2.0
    assert academic.styles['Title'].font.size.pt == 12
    assert academic.styles['Heading 1'].font.size.pt == 12
    assert 'Times New Roman' in academic.styles['Title'].element.xml
    assert 'asciiTheme' not in academic.styles['Title'].element.xml
    assert academic.styles['Heading 2'].font.italic is False
    assert round(academic.sections[0].left_margin.inches, 2) == 1.0

    resume = Document(BytesIO(generator.docx('# Jordan Ellis\n\n## Experience\n\n- Led product discovery.', profile='resume')))
    assert resume.styles['Normal'].font.name == 'Aptos'
    assert resume.sections[0].header.paragraphs[0].text == ''
    assert resume.sections[0].left_margin.inches < 0.7
    assert 'w:pBdr' not in resume.styles['Title'].element.xml


def test_explicit_resume_purpose_survives_a_fact_only_brief():
    brief = 'Led fraud operations for five years. Target role: product manager. Skills: SQL and analytics.'
    assert ArtifactService.profile_for('', brief, 'docx') == 'business'
    assert ArtifactService.profile_for('', brief, 'docx', 'resume') == 'resume'
    assert ArtifactService.normalize_purpose('RESUME') == 'resume'
    assert ArtifactService.normalize_purpose('arbitrary-client-value') is None
    guidance = ArtifactService.creation_guidance('docx', brief, 'resume')
    assert 'ATS-friendly' in guidance
    assert 'Quantify impact only when the user supplied the number' in guidance


def test_explicit_resume_purpose_controls_the_packaged_word_profile():
    files = DummyFiles()
    generator = ArtifactService(files)
    result = asyncio.run(generator.create(
        user_id='fixture-user',
        markdown='# Jordan Ellis\n\n## Experience\n\n- Led fraud operations.',
        format_name='docx',
        chat_id='fixture-chat',
        message_id='fixture-message',
        brief='Five years in operations. Target role: product manager. Skills: SQL.',
        purpose='resume',
    ))
    assert result['profile'] == 'resume'
    assert files.stored is not None
    assert files.stored['metadata']['profile'] == 'resume'
    document = Document(BytesIO(files.stored['data']))
    assert document.styles['Normal'].font.name == 'Aptos'
    assert document.sections[0].left_margin.inches < 0.7


def test_academic_references_use_a_standard_hanging_indent():
    academic = Document(BytesIO(service().docx(
        '# Evidence and Judgment\n\n## References\n\nCrump, G. (2026). Professional documents in practice.',
        profile='academic',
    )))
    reference = next(paragraph for paragraph in academic.paragraphs if paragraph.text.startswith('Crump, G.'))
    assert round(reference.paragraph_format.left_indent.inches, 2) == 0.5
    assert round(reference.paragraph_format.first_line_indent.inches, 2) == -0.5


def test_long_academic_title_is_complete_not_repeated_and_has_neutral_footer():
    title = 'Governance Before Scale: A Human-Centered Standard for Generative AI in Higher Education'
    academic = Document(BytesIO(service().docx(f'# {title}\n\nA finished paragraph.', profile='academic')))
    assert academic.paragraphs[0].text == title
    assert sum(paragraph.text == title for paragraph in academic.paragraphs) == 1
    footer = academic.sections[0].footer
    assert 'Ask Crump' not in footer.paragraphs[0].text
    assert 'PAGE' in footer._element.xml


def test_presentation_export_uses_custom_widescreen_layout_and_readable_type():
    deck = Presentation(BytesIO(service().pptx(SAMPLE)))
    assert round(deck.slide_width / deck.slide_height, 2) == round(16 / 9, 2)
    assert len(deck.slides) >= 3
    assert len(deck.slides[0].placeholders) == 0
    text_shapes = [shape for shape in deck.slides[0].shapes if getattr(shape, 'has_text_frame', False)]
    title_shape = next(shape for shape in text_shapes if shape.text == 'Product Brief')
    assert title_shape.text_frame.paragraphs[0].font.size.pt >= 50
    visible_text = '\n'.join(
        shape.text for slide in deck.slides for shape in slide.shapes
        if getattr(shape, 'has_text_frame', False)
    )
    assert 'ASK CRUMP' not in visible_text
    purpose = next(
        shape for shape in deck.slides[1].shapes
        if getattr(shape, 'has_text_frame', False) and shape.text == 'Purpose'
    )
    assert purpose.text_frame.paragraphs[0].font.size.pt >= 35


def test_presentation_export_has_color_rhythm_and_editable_native_charts():
    generator = service()
    rhythm = Presentation(BytesIO(generator.pptx(RHYTHM_DECK)))
    backgrounds = {
        str(slide.background.fill.fore_color.rgb)
        for slide in rhythm.slides
    }
    assert '171B24' in backgrounds
    assert 'FAF9F6' in backgrounds
    assert '202532' in backgrounds

    accent_fills = set()
    for slide in rhythm.slides:
        for shape in slide.shapes:
            try:
                value = shape.fill.fore_color.rgb
            except (AttributeError, TypeError, ValueError):
                continue
            if value:
                accent_fills.add(str(value))
    assert {'C9A95E', '1F6A67', 'B85C3B'} <= accent_fills

    chart_bytes = generator.pptx(NUMERIC_DECK)
    chart_deck = Presentation(BytesIO(chart_bytes))
    charts = [
        shape.chart for slide in chart_deck.slides for shape in slide.shapes
        if getattr(shape, 'has_chart', False)
    ]
    assert len(charts) == 1
    assert len(charts[0].series) == 2
    assert list(charts[0].plots[0].categories)[0].label == 'Organic'
    with zipfile.ZipFile(BytesIO(chart_bytes)) as package:
        chart_xml = package.read('ppt/charts/chart1.xml')
    assert b'<c:axId val="-' not in chart_xml
    assert b'<c:crossAx val="-' not in chart_xml


def test_presentation_uses_an_asymmetric_split_for_two_point_slides():
    deck = Presentation(BytesIO(service().pptx(RHYTHM_DECK)))
    first_content = deck.slides[1]
    dark_panels = [
        shape for shape in first_content.shapes
        if getattr(shape, 'shape_type', None) == 1
        and getattr(getattr(shape, 'fill', None), 'fore_color', None) is not None
        and str(shape.fill.fore_color.rgb) == '202532'
        and shape.width.inches > 5
    ]
    assert dark_panels
    assert dark_panels[0].left.inches > 6.5


def test_presentation_keeps_a_single_table_lead_with_its_chart():
    markdown = """# Activation Brief

## Activation is improving across channels

Referral is the strongest current acquisition path.

| Channel | Activation | Retention |
| --- | --- | --- |
| Organic | 42% | 31% |
| Referral | 58% | 46% |
| Paid | 36% | 24% |
| Partner | 49% | 39% |
"""
    deck = Presentation(BytesIO(service().pptx(markdown)))
    assert len(deck.slides) == 2
    data_slide = deck.slides[1]
    visible_text = '\n'.join(
        shape.text for shape in data_slide.shapes
        if getattr(shape, 'has_text_frame', False)
    )
    assert 'Activation is improving across channels' in visible_text
    assert 'Referral is the strongest current acquisition path.' in visible_text
    assert sum(1 for shape in data_slide.shapes if getattr(shape, 'has_chart', False)) == 1


def test_spreadsheet_export_is_structured_typed_filterable_and_safe():
    generator = service()
    workbook = load_workbook(BytesIO(generator.xlsx(SAMPLE)), data_only=False)
    assert workbook.sheetnames[0] == 'Overview'
    data_sheet = workbook[workbook.sheetnames[1]]
    assert workbook['Overview'].print_area == "'Overview'!$A$1:$C$10"
    assert workbook['Overview']['A3'].value == 'WORKBOOK INDEX'
    assert 'ASK CRUMP' not in ' '.join(str(cell.value or '') for row in workbook['Overview'].iter_rows() for cell in row)
    assert data_sheet.freeze_panes == 'A2'
    assert data_sheet.print_area
    assert data_sheet.sheet_view.showGridLines is False
    assert data_sheet.tables
    assert 'ASK CRUMP' not in (data_sheet.oddFooter.center.text or '')
    assert len(data_sheet.conditional_formatting) == 0
    assert data_sheet['B3'].value == 0.95
    assert data_sheet['B3'].number_format == '0.0%'
    assert data_sheet['B4'].value == 12500
    assert '$' in data_sheet['B4'].number_format
    assert isinstance(data_sheet['B5'].value, __import__('datetime').datetime)
    assert generator._typed_cell('=SUM(B2:B4)')[0] == '=SUM(B2:B4)'
    assert generator._typed_cell('=HYPERLINK("https://bad.example","open")')[0].startswith("'=")
    assert generator._typed_cell("='[external.xlsx]Sheet1'!A1")[0].startswith("'=")


def test_pdf_export_is_extractable_and_paginated():
    reader = PdfReader(BytesIO(service().pdf(SAMPLE)))
    assert reader.metadata.author == 'Ask Crump'
    extracted = '\n'.join(page.extract_text() or '' for page in reader.pages)
    assert 'Product Brief' in extracted
    assert 'ASK CRUMP' not in extracted


def test_creation_guidance_guards_truth_and_matches_the_format():
    resume = ArtifactService.creation_guidance('docx', 'Create my résumé')
    academic = ArtifactService.creation_guidance('docx', 'Write a college essay with APA citations')
    slides = ArtifactService.creation_guidance('pptx', 'Create a board presentation')
    workbook = ArtifactService.creation_guidance('xlsx', 'Build a financial model')
    assert 'Never invent citations' in resume
    assert 'ATS-friendly' in resume and 'user supplied the number' in resume
    assert 'Cite only sources' in academic and 'defensible thesis' in academic
    assert 'one clear takeaway headline per slide' in slides
    assert 'editable chart' in slides
    assert 'Never invent business or financial data' in workbook
